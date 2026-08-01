"""One-off generator for gibberish sample PDF/PPTX files used to test file upload + preview."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

OUT = Path(__file__).parent

PDF_DECKS = [
    ("sample-deck-quantum-toaster.pdf", "Quantum Toaster", [
        "Quantum Toaster — Pitch Deck",
        "Problem: bread is toasted using outdated classical physics.",
        "Solution: quantum tunneling toasts bread before you even want it.",
        "Market size: every kitchen, theoretically, in every universe.",
        "Ask: $2M seed round, paid in theoretical toast.",
    ]),
    ("sample-deck-ferret-as-a-service.pdf", "Ferret-as-a-Service", [
        "Ferret-as-a-Service — Pitch Deck",
        "Problem: emotional support is expensive and illiquid.",
        "Solution: rent a ferret by the hour, cloud-dispatched.",
        "Traction: 14 ferrets, 3 mildly satisfied customers.",
        "Ask: $500K to expand the ferret fleet.",
    ]),
]

PPTX_DECKS = [
    ("sample-slides-echomood.pptx", "EchoMood", [
        "EchoMood",
        "We detect mood from sighs.",
        "Thank you. Questions? (please sigh into the mic)",
    ]),
    ("sample-slides-noodlenet.pptx", "NoodleNet", [
        "NoodleNet",
        "Peer-to-peer noodle recipes, ranked by spice tolerance.",
        "Thank you. May your noodles be ever slurpable.",
    ]),
]


def make_pdf(filename: str, title: str, lines: list[str]) -> None:
    path = OUT / filename
    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 18)
    c.drawString(72, height - 100, title)
    c.setFont("Helvetica", 12)
    y = height - 140
    for line in lines:
        c.drawString(72, y, line)
        y -= 24
    c.showPage()
    c.save()
    print(f"wrote {path}")


def make_pptx(filename: str, title: str, slide_texts: list[str]) -> None:
    path = OUT / filename
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for i, text in enumerate(slide_texts):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title if i == 0 else text.split(".")[0]
        body = slide.placeholders[1].text_frame
        body.text = text
    prs.save(str(path))
    print(f"wrote {path}")


if __name__ == "__main__":
    for filename, title, lines in PDF_DECKS:
        make_pdf(filename, title, lines)
    for filename, title, texts in PPTX_DECKS:
        make_pptx(filename, title, texts)
